"""
Stage 06 — Render (deterministic, python-pptx).

Uses the GDT DECK.pptx template as the base so every generated deck
inherits the Unilever slide master, fonts, backgrounds, and decorations.
"""
import asyncio
import base64
import io
import json
import re
import zipfile
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from models.schemas import DeckSpec, SlideSpec, VisualSpec
from brand.unilever_2026 import BRAND, TEMPLATE_LAYOUT_MAP

TEMPLATE_PATH = Path(__file__).parent.parent / "template" / "GDT DECK.pptx"

# ── Brand colour helpers ──────────────────────────────────────────────────────
def _rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

CHART_PALETTE = [
    _rgb(BRAND["primary"]),    # #0066CC
    _rgb(BRAND["dark"]),       # #133061
    _rgb(BRAND["purple"]),     # #8651DF
    _rgb(BRAND["teal"]),       # #008090
    _rgb(BRAND["green"]),      # #2B911C
    _rgb(BRAND["orange"]),     # #DA5700
]


# ── Public entry point ────────────────────────────────────────────────────────

async def render_pptx(deck: DeckSpec, job_id: str, upload_dir: str) -> Path:
    out_path = Path(upload_dir) / f"{job_id}.pptx"
    loop = asyncio.get_event_loop()
    await loop.run_in_executor(None, _render_sync, deck, out_path)
    return out_path


def _render_sync(deck: DeckSpec, out_path: Path) -> None:
    prs = Presentation(str(TEMPLATE_PATH))
    _clear_slides(prs)

    total = len(deck.slides)
    for spec in deck.slides:
        _add_slide(prs, spec, total)

    out_path.parent.mkdir(parents=True, exist_ok=True)

    # Save to a buffer, then rewrite the ZIP removing orphaned template slide
    # files — they cause PowerPoint's "repair" warning on open.
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    _clean_pptx_zip(buf, out_path)

    # Sidecar: deck title used by the download endpoint for Content-Disposition
    meta = out_path.with_suffix(".meta.json")
    meta.write_text(
        json.dumps({"title": deck.title or "Unilever Deck"}, ensure_ascii=False),
        encoding="utf-8",
    )


def _clean_pptx_zip(src: io.BytesIO, dest: Path) -> None:
    """
    Rewrite the PPTX ZIP so it contains only the slides that are actually in
    the deck's sldIdLst.  This removes the 25 example slides that were in the
    GDT DECK template and caused PowerPoint's "repair" warning.

    Strategy:
    1. Build a content dict (last-write-wins for duplicates, e.g. slide1.xml
       that appears twice because python-pptx reuses numbers).
    2. Read the sldIdLst from presentation.xml to find which rIds are active.
    3. Look up those rIds in the rels file to find the actual slide filenames.
    4. Strip any slide*.xml / slide*.xml.rels files not in that wanted set.
    5. Rewrite presentation.xml.rels removing relationships for old slides
       so there are no dangling rels entries pointing at missing files.
    """
    with zipfile.ZipFile(src, "r") as zin:
        # ── build content dict: last entry wins for duplicate filenames ───────
        content: dict[str, tuple] = {}
        for item in zin.infolist():
            content[item.filename] = (item, zin.read(item.filename))

        prs_xml  = content["ppt/presentation.xml"][1].decode("utf-8")
        rels_xml = content["ppt/_rels/presentation.xml.rels"][1].decode("utf-8")

        # ── active rIds come from sldIdLst, NOT from rels ─────────────────────
        active_rids = set(re.findall(
            r'<p:sldId\b[^>]+\br:id="([^"]+)"', prs_xml
        ))

        # ── map rId → target path ─────────────────────────────────────────────
        rid_to_target: dict[str, str] = {}
        for m in re.finditer(
            r'<Relationship[^>]+\bId="([^"]+)"[^>]+\bTarget="([^"]+)"', rels_xml
        ):
            rid_to_target[m.group(1)] = m.group(2)

        # ── filenames we want to keep ─────────────────────────────────────────
        wanted: set[str] = set()
        for rid in active_rids:
            if rid not in rid_to_target:
                continue
            target = rid_to_target[rid]          # e.g. "slides/slide26.xml"
            full   = f"ppt/{target}"             # "ppt/slides/slide26.xml"
            wanted.add(full)
            fname    = target.rsplit("/", 1)[-1]  # "slide26.xml"
            slide_dir = target.rsplit("/", 1)[0]  # "slides"
            wanted.add(f"ppt/{slide_dir}/_rels/{fname}.rels")

        def _keep(name: str) -> bool:
            if re.match(r"ppt/slides/slide\d+\.xml$", name):
                return name in wanted
            if re.match(r"ppt/slides/_rels/slide\d+\.xml\.rels$", name):
                return name in wanted
            return True

        # ── clean the rels file: strip stale slide relationships ─────────────
        def _clean_rels(xml: str) -> str:
            def _sub(m: re.Match) -> str:
                tag = m.group(0)
                rid_m = re.search(r'\bId="([^"]+)"', tag)
                if rid_m and rid_m.group(1) not in active_rids:
                    if "slides/slide" in tag:
                        return ""
                return tag
            return re.sub(r'<Relationship\b[^/]*/>', _sub, xml)

        cleaned_rels = _clean_rels(rels_xml)

        # ── write clean ZIP ───────────────────────────────────────────────────
        with zipfile.ZipFile(str(dest), "w", zipfile.ZIP_DEFLATED) as zout:
            for name, (item, data) in content.items():
                if not _keep(name):
                    continue
                if name == "ppt/_rels/presentation.xml.rels":
                    zout.writestr(item, cleaned_rels.encode("utf-8"))
                else:
                    zout.writestr(item, data)


def _clear_slides(prs: Presentation) -> None:
    """Remove example slides from the template, preserving master/theme."""
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        sld_id_lst.remove(sld_id)


# ── Per-slide rendering ───────────────────────────────────────────────────────

def _add_slide(prs: Presentation, spec: SlideSpec, total: int) -> None:
    layout_idx = TEMPLATE_LAYOUT_MAP.get(spec.layout, 7)
    layout = prs.slide_masters[0].slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)

    headline = spec.copy.get("headline", "")
    body     = spec.copy.get("body", "")
    section  = spec.section.upper()

    # ── Cover (Layout 3: '1: Cover') ─────────────────────────────────────────
    # Reference style: big title + date "Jun 2026" — no subtitle
    if spec.layout == "cover":
        # Optional AI hero image as a full-bleed background (sent behind text).
        img = _decode_image(spec.visual)
        if img:
            _add_background_image(slide, img)
        _set_ph(slide, 0, headline)
        _hide_ph(slide, 1)
        today = date.today()
        _set_ph(slide, 14, today.strftime("%b %Y"))

    # ── Agenda (Section - Sky, Layout 5) ─────────────────────────────────────
    elif spec.layout == "agenda":
        _set_ph(slide, 1, section)
        _set_ph(slide, 0, headline)
        if body:
            items = [b.strip() for b in body.replace("·", "\n").split("\n") if b.strip()]
            _add_textbox(slide,
                         "\n".join(f"▸  {item}" for item in items[:8]),
                         Inches(0.54), Inches(0.6), Inches(8.5), Inches(2.1),
                         size=Pt(13))

    # ── Close (Section - Earth, Layout 6) ────────────────────────────────────
    elif spec.layout == "close":
        _set_ph(slide, 1, section)
        _set_ph(slide, 0, headline)
        if body:
            lines = [b.strip() for b in body.replace("·", "\n").split("\n") if b.strip()]
            _add_textbox(slide,
                         "\n".join(f"→  {l}" for l in lines[:4]),
                         Inches(0.54), Inches(0.6), Inches(9.0), Inches(2.0),
                         size=Pt(13))

    # ── Split-left (Layout 8: '3: Left') ─────────────────────────────────────
    # Title left at (0.406", 1.311") size 5.51"x1.447" → bottom edge 2.758"
    # Body text: left column at 3.0" (safe gap below title)
    # Chart: right column starting at 1.4" for maximum height
    # idx=13 OBJECT placeholder (dashed box) must be hidden if unused
    elif spec.layout == "left":
        _set_ph(slide, 15, section)
        _set_title(slide, headline)             # capped at 26pt — stays in left column
        _hide_ph(slide, 13)                     # hide "Type body copy here" dashed box
        has_vis = spec.visual and _has_chart_data(spec.visual)
        if has_vis:
            if body:
                _add_textbox(slide, body,
                             Inches(0.41), Inches(3.0), Inches(5.1), Inches(3.9),
                             size=Pt(12))
            _draw_chart(slide, spec.visual,
                        Inches(5.5), Inches(1.4), Inches(7.5), Inches(5.7))
        else:
            if body:
                _add_textbox(slide, body,
                             Inches(0.41), Inches(3.0), Inches(12.3), Inches(3.9),
                             size=Pt(13))

    # ── Split-right (Layout 9: '3: Right') ───────────────────────────────────
    # Title right at (6.573", 1.533") size 5.603"x1.447" → bottom edge 2.980"
    # Chart: left column (independent of right-side title)
    # Text: right column below title at 3.1"
    elif spec.layout == "right":
        _set_ph(slide, 15, section)
        _set_title(slide, headline)             # capped at 26pt — stays in right column
        _hide_ph(slide, 13)                     # hide dashed OBJECT placeholder
        has_vis = spec.visual and _has_chart_data(spec.visual)
        if has_vis:
            _draw_chart(slide, spec.visual,
                        Inches(0.41), Inches(1.4), Inches(5.8), Inches(5.7))
            if body:
                _add_textbox(slide, body,
                             Inches(6.4), Inches(3.1), Inches(6.5), Inches(3.8),
                             size=Pt(12))
        else:
            if body:
                _add_textbox(slide, body,
                             Inches(0.41), Inches(3.1), Inches(12.3), Inches(3.8),
                             size=Pt(13))

    # ── Text slides (Layout 7: 3: Title) ─────────────────────────────────────
    # Title at (0.406", 0.932") h=0.689" — at 26pt wraps to ≤2 lines ≈ 2.3" bottom
    # Safe body start: 3.0" — guarantees no collision with 2-line title
    elif spec.layout in ("context", "findings", "recommendation"):
        _set_ph(slide, 15, section)
        _set_title(slide, headline)
        img = _decode_image(spec.visual)
        if img:
            # Body on the left, generated image on the right.
            if body:
                _add_textbox(slide, body,
                             Inches(0.41), Inches(3.0), Inches(7.2), Inches(3.9),
                             size=Pt(13))
            _add_picture_fit(slide, img, Inches(8.0), Inches(2.7), Inches(4.8), Inches(3.6))
        elif body:
            _add_textbox(slide, body,
                         Inches(0.41), Inches(3.0), Inches(12.3), Inches(3.9),
                         size=Pt(13))

    # ── Chart slide (Layout 7) ────────────────────────────────────────────────
    elif spec.layout == "chart":
        _set_ph(slide, 15, section)
        _set_title(slide, headline)
        cx, cy, cw, ch = Inches(0.5), Inches(1.9), Inches(12.2), Inches(5.0)
        if spec.visual and _has_chart_data(spec.visual):
            _draw_chart(slide, spec.visual, cx, cy, cw, ch)
        else:
            caption = spec.visual.caption if spec.visual else ""
            _add_textbox(slide, f"[Chart: {caption}]",
                         cx, cy + ch / 2, cw, Inches(0.5),
                         size=Pt(12), color=BRAND["muted"])

    # ── Metrics slide ─────────────────────────────────────────────────────────
    elif spec.layout == "metrics":
        _set_ph(slide, 15, section)
        _set_title(slide, headline)
        _draw_metrics(slide, body)

    # ── Diagram slide ─────────────────────────────────────────────────────────
    elif spec.layout == "diagram":
        _set_ph(slide, 15, section)
        _set_title(slide, headline)
        if spec.visual and spec.visual.data and spec.visual.data.get("nodes"):
            _draw_flow_diagram(slide, spec.visual,
                               Inches(0.5), Inches(2.0), Inches(12.2), Inches(4.7))
        else:
            _add_textbox(slide, "[Diagram: generating]",
                         Inches(0.5), Inches(3.0), Inches(12.2), Inches(0.5),
                         size=Pt(12), color=BRAND["muted"])

    # ── Roadmap slide ─────────────────────────────────────────────────────────
    elif spec.layout == "roadmap":
        _set_ph(slide, 15, section)
        _set_title(slide, headline)
        _draw_roadmap(slide, body)

    # ── Achievement slide (content-fidelity) ───────────────────────────────────
    elif spec.layout == "achievement":
        _set_ph(slide, 15, section)
        _set_title(slide, headline)
        _draw_achievement(slide, spec.copy)

    # ── Fallback ──────────────────────────────────────────────────────────────
    else:
        _set_ph(slide, 15, section)
        _set_title(slide, headline)
        if body:
            _add_textbox(slide, body,
                         Inches(0.41), Inches(3.0), Inches(12.3), Inches(3.9),
                         size=Pt(13))

    # ── Speaker notes (all layouts) ──────────────────────────────────────────
    if spec.speaker_notes:
        _set_notes(slide, spec.speaker_notes)


# ── Chart rendering ───────────────────────────────────────────────────────────

def _has_chart_data(visual: VisualSpec) -> bool:
    d = visual.data or {}
    if visual.type == "table":
        return bool(d.get("headers") or d.get("rows"))
    return bool(d.get("labels") and d.get("series"))


def _draw_chart(slide, visual: VisualSpec, x, y, w, h) -> None:
    if visual.type == "table":
        _draw_table(slide, visual, x, y, w, h)
        return
    try:
        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

        TYPE_MAP = {
            "bar":         XL_CHART_TYPE.BAR_CLUSTERED,
            "column":      XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line":        XL_CHART_TYPE.LINE_MARKERS,
            "pie":         XL_CHART_TYPE.PIE,
            "donut":       XL_CHART_TYPE.DOUGHNUT,
            "donut_chart": XL_CHART_TYPE.DOUGHNUT,
        }
        chart_type = TYPE_MAP.get(visual.type, XL_CHART_TYPE.COLUMN_CLUSTERED)

        d = visual.data
        labels      = d.get("labels", [])
        series_list = d.get("series", [])

        cd = ChartData()
        cd.categories = labels
        for s in series_list:
            vals = [float(v) if v is not None else 0.0 for v in s.get("values", [])]
            cd.add_series(s.get("name", ""), vals)

        chart_h = h - Inches(0.4)
        chart_frame = slide.shapes.add_chart(chart_type, x, y, w, chart_h, cd)
        chart = chart_frame.chart
        chart.has_title = False
        chart.has_legend = (len(series_list) > 1 or
                            chart_type in (XL_CHART_TYPE.PIE, XL_CHART_TYPE.DOUGHNUT))
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False

        _color_chart_series(chart, chart_type)

        # Caption + provenance: cite the source when the data is grounded.
        caption = visual.caption
        if visual.grounded and visual.source_ref:
            caption = f"{caption}  ·  Source: {visual.source_ref}".strip(" ·")
        if caption:
            _add_textbox(slide, caption,
                         x, y + chart_h + Inches(0.05), w, Inches(0.3),
                         size=Pt(9), color=BRAND["muted"])

    except Exception:
        _add_textbox(slide, visual.caption or "[Chart]",
                     x + Inches(0.2), y + h / 2 - Inches(0.25),
                     w - Inches(0.4), Inches(0.5),
                     size=Pt(11), color=BRAND["muted"])


def _color_chart_series(chart, chart_type) -> None:
    """
    Apply distinct brand colours to chart data.

    Pie / Donut: each SLICE gets its own colour (one series, N points).
    Single-series bar/column: each BAR gets its own colour for visual variety.
    Multi-series: each series gets one colour (standard approach).
    """
    from pptx.enum.chart import XL_CHART_TYPE

    is_pie = chart_type in (XL_CHART_TYPE.PIE, XL_CHART_TYPE.DOUGHNUT)
    series_list = list(chart.series)

    if is_pie and series_list:
        # Colour each slice independently
        for i, point in enumerate(series_list[0].points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = CHART_PALETTE[i % len(CHART_PALETTE)]
    elif len(series_list) == 1:
        # Single-series bar/column — colour each bar differently
        try:
            for i, point in enumerate(series_list[0].points):
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = CHART_PALETTE[i % len(CHART_PALETTE)]
        except Exception:
            # Fall back to series-level colour if point API not available
            series_list[0].format.fill.solid()
            series_list[0].format.fill.fore_color.rgb = CHART_PALETTE[0]
    else:
        # Multi-series: one colour per series
        for i, series in enumerate(series_list):
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = CHART_PALETTE[i % len(CHART_PALETTE)]


def _draw_table(slide, visual: VisualSpec, x, y, w, h) -> None:
    try:
        d = visual.data or {}
        headers = d.get("headers", [])
        rows    = d.get("rows", [])
        if not headers and not rows:
            return

        n_cols = len(headers) if headers else (len(rows[0]) if rows else 1)
        n_rows = min(len(rows) + (1 if headers else 0), 12)

        tbl = slide.shapes.add_table(n_rows, n_cols, x, y, w, h).table

        if headers:
            for j, hdr in enumerate(headers[:n_cols]):
                cell = tbl.cell(0, j)
                cell.text = str(hdr)
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(BRAND["primary"])
                p = cell.text_frame.paragraphs[0]
                run = p.runs[0] if p.runs else p.add_run()
                run.font.color.rgb = _rgb(BRAND["white"])
                run.font.bold = True
                run.font.size = Pt(11)
                run.font.name = BRAND["font_body"]

        alt = (_rgb(BRAND["surface"]), _rgb(BRAND["white"]))
        for i, row in enumerate(rows[:n_rows]):
            ri = i + (1 if headers else 0)
            if ri >= n_rows:
                break
            for j, val in enumerate(row[:n_cols]):
                cell = tbl.cell(ri, j)
                cell.text = str(val)
                cell.fill.solid()
                cell.fill.fore_color.rgb = alt[i % 2]
                p = cell.text_frame.paragraphs[0]
                run = p.runs[0] if p.runs else p.add_run()
                run.font.size = Pt(10)
                run.font.name = BRAND["font_body"]
                run.font.color.rgb = _rgb(BRAND["dark"])

    except Exception:
        _add_textbox(slide, visual.caption or "[Table]",
                     x + Inches(0.2), y + h / 2,
                     w, Inches(0.4), size=Pt(11), color=BRAND["muted"])


# ── Flow diagram ──────────────────────────────────────────────────────────────

def _draw_flow_diagram(slide, visual: VisualSpec, x, y, w, h) -> None:
    try:
        nodes = (visual.data or {}).get("nodes", [])[:6]
        n = len(nodes)
        if not n:
            return

        node_h  = Inches(0.9)
        arrow_w = Inches(0.22)
        node_w  = (w - arrow_w * (n - 1)) / n
        center_y = y + h / 2 - node_h / 2

        for i, label in enumerate(nodes):
            nx = x + i * (node_w + arrow_w)
            col = CHART_PALETTE[i % len(CHART_PALETTE)]

            box = slide.shapes.add_shape(1, nx, center_y, node_w, node_h)
            box.fill.solid()
            box.fill.fore_color.rgb = col
            box.line.fill.background()
            box.shadow.inherit = False

            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = label
            run.font.name = BRAND["font_body"]
            run.font.size = Pt(11)
            run.font.bold = True
            run.font.color.rgb = _rgb(BRAND["white"])

            if i < n - 1:
                ax = nx + node_w
                ay = center_y + node_h / 2 - Inches(0.03)
                arr = slide.shapes.add_shape(1, ax, ay, arrow_w, Inches(0.06))
                arr.fill.solid()
                arr.fill.fore_color.rgb = _rgb(BRAND["border"])
                arr.line.fill.background()
                arr.shadow.inherit = False

        if visual.caption:
            cap_y = center_y + node_h + Inches(0.12)
            _add_textbox(slide, visual.caption, x, cap_y, w, Inches(0.3),
                         size=Pt(9), color=BRAND["muted"])

    except Exception:
        _add_textbox(slide, "[Diagram]", x, y + h / 2, w, Inches(0.4),
                     size=Pt(11), color=BRAND["muted"])


# ── Achievement blocks (Overview / What went well / What could be better) ──────

def _tint(hex_str: str, ratio: float = 0.86) -> RGBColor:
    """Lighten a brand colour toward white for a soft block background."""
    h = hex_str.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return RGBColor(int(r + (255 - r) * ratio),
                    int(g + (255 - g) * ratio),
                    int(b + (255 - b) * ratio))


def _draw_block(slide, header: str, body: str, x, y, w, h, accent_hex: str) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = _tint(accent_hex)
    box.line.color.rgb = _rgb(accent_hex)
    box.line.width = Pt(1)
    box.shadow.inherit = False

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.22)
    tf.margin_right = Inches(0.22)
    tf.margin_top = Inches(0.16)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = header
    run.font.bold = True
    run.font.size = Pt(11)
    run.font.name = BRAND["font_body"]
    run.font.color.rgb = _rgb(accent_hex)

    if body:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(6)
        r2 = p2.add_run()
        r2.text = body
        r2.font.size = Pt(11)
        r2.font.name = BRAND["font_body"]
        r2.font.color.rgb = _rgb(BRAND["dark"])


def _draw_achievement(slide, copy: dict) -> None:
    overview = (copy.get("overview") or "").strip()
    went = (copy.get("went_well") or "").strip()
    could = (copy.get("could_better") or "").strip()
    # If the model only produced a plain body, fall back to a simple text block.
    if not (overview or went or could):
        body = (copy.get("body") or "").strip()
        if body:
            _add_textbox(slide, body, Inches(0.5), Inches(1.9), Inches(12.3), Inches(4.8), size=Pt(13))
        return

    y = Inches(1.55)
    if overview:
        _add_textbox(slide, overview, Inches(0.5), y, Inches(12.3), Inches(1.25),
                     size=Pt(12), color=BRAND["dark"])
        y = Inches(3.0)
    else:
        y = Inches(1.9)

    box_h = Inches(3.7) if overview else Inches(4.6)
    if went and could:
        _draw_block(slide, "✓  WHAT WENT WELL", went, Inches(0.5), y, Inches(6.05), box_h, BRAND["green"])
        _draw_block(slide, "△  WHAT COULD BE BETTER", could, Inches(6.78), y, Inches(6.05), box_h, BRAND["orange"])
    elif went:
        _draw_block(slide, "✓  WHAT WENT WELL", went, Inches(0.5), y, Inches(12.3), box_h, BRAND["green"])
    elif could:
        _draw_block(slide, "△  WHAT COULD BE BETTER", could, Inches(0.5), y, Inches(12.3), box_h, BRAND["orange"])


# ── Metrics boxes ─────────────────────────────────────────────────────────────

def _draw_metrics(slide, body: str) -> None:
    parts = [m.strip() for m in body.split("·") if m.strip()]
    if not parts:
        return
    grid = [(0, 0), (1, 0), (0, 1), (1, 1)]
    for i, metric in enumerate(parts[:4]):
        col, row = grid[i]
        bx = Inches(0.5 + col * 6.3)
        by = Inches(1.88 + row * 2.55)
        box = slide.shapes.add_shape(1, bx, by, Inches(5.9), Inches(2.2))
        box.fill.solid()
        box.fill.fore_color.rgb = _rgb(BRAND["surface"])
        box.line.color.rgb = _rgb(BRAND["border"])
        box.line.width = 9525
        box.shadow.inherit = False

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = metric
        run.font.name = BRAND["font_body"]
        run.font.size = Pt(13)
        run.font.color.rgb = _rgb(BRAND["primary"])


# ── Roadmap phases ────────────────────────────────────────────────────────────

def _draw_roadmap(slide, body: str) -> None:
    phases = [p.strip() for p in body.split("·") if p.strip()]
    if not phases:
        return
    n       = min(len(phases), 6)
    gap     = Inches(0.15)
    total_w = Inches(12.3)
    box_w   = (total_w - gap * (n - 1)) / n
    sx      = Inches(0.5)

    for i, phase in enumerate(phases[:n]):
        px  = sx + i * (box_w + gap)
        col = CHART_PALETTE[i % len(CHART_PALETTE)]

        box = slide.shapes.add_shape(1, px, Inches(2.1), box_w, Inches(1.4))
        box.fill.solid()
        box.fill.fore_color.rgb = col
        box.line.fill.background()
        box.shadow.inherit = False

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = phase
        run.font.name = BRAND["font_body"]
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = _rgb(BRAND["white"])


# ── Helpers ────────────────────────────────────────────────────────────────────

def _set_ph(slide, idx: int, text: str) -> None:
    try:
        ph = slide.placeholders[idx]
        ph.text = text
    except (KeyError, IndexError):
        pass


def _set_title(slide, text: str, max_pt: int = 26) -> None:
    """
    Set the title placeholder (idx=0) and cap its font size so long titles
    stay within their bounding box and never collide with body content.
    At 26pt Century Gothic a 60-char title fits on one line in the 12" placeholder.
    """
    try:
        ph = slide.placeholders[0]
        ph.text = text
        tf = ph.text_frame
        # Apply size cap at the paragraph level (affects all runs without
        # creating per-run overrides that fight with the theme).
        for para in tf.paragraphs:
            para.font.size = Pt(max_pt)
    except (KeyError, IndexError, AttributeError):
        pass


def _decode_image(visual) -> io.BytesIO | None:
    """Return a BytesIO of an AI-generated image, or None."""
    if not visual or visual.type != "image":
        return None
    b64 = (visual.data or {}).get("image_b64")
    if not b64:
        return None
    try:
        return io.BytesIO(base64.b64decode(b64))
    except Exception:
        return None


def _add_background_image(slide, stream: io.BytesIO) -> None:
    """Add a full-bleed image and send it behind all other shapes (cover hero)."""
    try:
        pic = slide.shapes.add_picture(stream, 0, 0, Inches(13.333), Inches(7.5))
        spTree = slide.shapes._spTree
        spTree.remove(pic._element)
        spTree.insert(2, pic._element)  # behind placeholders (after nvGrpSpPr, grpSpPr)
    except Exception:
        pass


def _add_picture_fit(slide, stream: io.BytesIO, x, y, w, h) -> None:
    try:
        slide.shapes.add_picture(stream, x, y, w, h)
    except Exception:
        pass


def _set_notes(slide, text: str) -> None:
    """Write presenter script into the slide's notes pane."""
    try:
        notes = slide.notes_slide
        notes.notes_text_frame.text = text
    except Exception:
        pass


def _hide_ph(slide, idx: int) -> None:
    """Remove a placeholder element from the slide XML so it doesn't render."""
    try:
        ph = slide.placeholders[idx]
        sp = ph._element
        sp.getparent().remove(sp)
    except (KeyError, IndexError, AttributeError):
        pass


def _add_textbox(slide, text: str, x, y, w, h,
                 size=None,
                 color: str = BRAND["primary"],
                 bold: bool = False,
                 align=PP_ALIGN.LEFT) -> None:
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = BRAND["font_body"]
    if size:
        run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
