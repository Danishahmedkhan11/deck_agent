"""
/api/upload   — multipart file upload → returns file_id
/api/generate — POST body → streams SSE pipeline events → returns download URL
/api/download/{job_id} — download the rendered .pptx
"""
import asyncio
import json
import os
import re
import uuid
from pathlib import Path

import aiofiles
from fastapi import APIRouter, Body, File, HTTPException, UploadFile
from fastapi.responses import FileResponse
from sse_starlette.sse import EventSourceResponse

from agents.pipeline import run_pipeline, compute_outline, build_from_context
from models.config import Settings
from models.schemas import GenerateRequest, StageEvent, StageStatus

router = APIRouter(tags=["generate"])
settings = Settings()

# in-memory job store (swap for Redis in production)
_jobs: dict[str, dict] = {}


# ── file upload ──────────────────────────────────────────────────────────────

@router.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    max_bytes = settings.max_upload_mb * 1024 * 1024
    file_id = str(uuid.uuid4())
    ext = Path(file.filename or "upload").suffix
    dest = Path(settings.upload_dir) / f"{file_id}{ext}"
    dest.parent.mkdir(parents=True, exist_ok=True)

    size = 0
    async with aiofiles.open(dest, "wb") as f:
        while chunk := await file.read(65536):
            size += len(chunk)
            if size > max_bytes:
                dest.unlink(missing_ok=True)
                raise HTTPException(413, f"File exceeds {settings.max_upload_mb} MB limit")
            await f.write(chunk)

    return {"file_id": file_id, "filename": file.filename, "size": size}


# ── generate (SSE streaming) ─────────────────────────────────────────────────

@router.post("/generate")
async def generate(req: GenerateRequest):
    job_id = str(uuid.uuid4())
    _jobs[job_id] = {"status": "running", "download_url": ""}

    async def event_stream():
        try:
            async for event in run_pipeline(req, job_id, settings):
                data = event.model_dump_json()
                yield {"event": "stage", "data": data}

            out_path = Path(settings.upload_dir) / f"{job_id}.pptx"
            if out_path.exists():
                _jobs[job_id] = {
                    "status": "done",
                    "download_url": f"/api/download/{job_id}",
                }
            done_event = StageEvent(
                stage="complete",
                status=StageStatus.done,
                message=f"/api/download/{job_id}",
                progress=100,
            )
            yield {"event": "done", "data": done_event.model_dump_json()}

        except Exception as exc:
            err = StageEvent(stage="error", status=StageStatus.error, message=str(exc))
            yield {"event": "error", "data": err.model_dump_json()}

    return EventSourceResponse(event_stream())


# ── outline review (two-phase) ────────────────────────────────────────────────

@router.post("/outline")
async def outline(req: GenerateRequest):
    """Phase 1: plan the deck and return an editable outline. Persists the
    phase-1 context so /api/build can finish the deck from the edited outline."""
    job_id = str(uuid.uuid4())
    context = await compute_outline(req, settings)

    # persist context for the build phase
    ctx_path = Path(settings.upload_dir) / f"{job_id}.context.json"
    ctx_path.parent.mkdir(parents=True, exist_ok=True)
    ctx_path.write_text(json.dumps(context), encoding="utf-8")

    deck = context["deck_spec"]
    outline_slides = [
        {
            "orig_index": s.get("index", i),
            "section": s.get("section", ""),
            "layout": s.get("layout", "findings"),
            "headline": (s.get("copy") or {}).get("headline", ""),
            "body": (s.get("copy") or {}).get("body", ""),
        }
        for i, s in enumerate(deck.get("slides", []))
    ]
    return {
        "job_id": job_id,
        "title": deck.get("title", ""),
        "intent": context.get("intent", {}),
        "auto_count": context.get("auto_count", False),
        "slides": outline_slides,
    }


@router.post("/build")
async def build(payload: dict = Body(...)):
    """Phase 2: build the full deck from a (possibly edited) outline."""
    job_id = payload.get("job_id")
    edited = payload.get("slides") or []
    if not job_id:
        raise HTTPException(400, "job_id is required")

    ctx_path = Path(settings.upload_dir) / f"{job_id}.context.json"
    if not ctx_path.exists():
        raise HTTPException(404, "Outline expired or not found — regenerate the outline")
    context = json.loads(ctx_path.read_text(encoding="utf-8"))

    # Apply the user's edits: reorder / drop / rewrite headlines, keeping the
    # original slide's layout/section/visual unless the edit overrides them.
    orig = {s.get("index"): s for s in context["deck_spec"].get("slides", [])}
    new_slides = []
    for pos, item in enumerate(edited):
        base = dict(orig.get(item.get("orig_index"), {}))
        base["index"] = pos
        base["section"] = item.get("section", base.get("section", ""))
        base["layout"] = item.get("layout", base.get("layout", "findings"))
        copy = dict(base.get("copy") or {})
        copy["headline"] = item.get("headline", copy.get("headline", ""))
        if item.get("body") is not None:
            copy["body"] = item.get("body")
        base["copy"] = copy
        new_slides.append(base)
    if new_slides:
        context["deck_spec"]["slides"] = new_slides
        context["deck_spec"]["slide_count"] = len(new_slides)

    _jobs[job_id] = {"status": "running", "download_url": ""}

    async def event_stream():
        try:
            async for event in build_from_context(context, job_id, settings):
                yield {"event": "stage", "data": event.model_dump_json()}
            out_path = Path(settings.upload_dir) / f"{job_id}.pptx"
            if out_path.exists():
                _jobs[job_id] = {"status": "done", "download_url": f"/api/download/{job_id}"}
            done_event = StageEvent(stage="complete", status=StageStatus.done,
                                    message=f"/api/download/{job_id}", progress=100)
            yield {"event": "done", "data": done_event.model_dump_json()}
        except Exception as exc:
            err = StageEvent(stage="error", status=StageStatus.error, message=str(exc))
            yield {"event": "error", "data": err.model_dump_json()}

    return EventSourceResponse(event_stream())


# ── download ─────────────────────────────────────────────────────────────────

@router.get("/download/{job_id}")
async def download(job_id: str):
    path = Path(settings.upload_dir) / f"{job_id}.pptx"
    if not path.exists():
        raise HTTPException(404, "Deck not found or still generating")

    # Read title from sidecar written by render stage
    filename = "Unilever Deck.pptx"
    meta_path = path.with_suffix(".meta.json")
    if meta_path.exists():
        try:
            meta = json.loads(meta_path.read_text(encoding="utf-8"))
            raw = meta.get("title", "Unilever Deck")
            # Sanitise: keep letters, digits, spaces, hyphens
            slug = re.sub(r"[^\w\s\-]", "", raw).strip()
            slug = re.sub(r"\s+", " ", slug)[:80]
            if slug:
                filename = f"{slug}.pptx"
        except Exception:
            pass

    return FileResponse(
        path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=filename,
    )


@router.get("/jobs/{job_id}")
async def job_status(job_id: str):
    if job_id not in _jobs:
        raise HTTPException(404, "Job not found")
    return _jobs[job_id]


# ── preview (slide data for in-browser preview) ───────────────────────────────

@router.get("/preview/{job_id}")
async def preview(job_id: str):
    # Preferred: serve from the persisted DeckSpec — faithful to what was planned,
    # and carries intent + quality score. Falls back to parsing the .pptx.
    spec_path = Path(settings.upload_dir) / f"{job_id}.spec.json"
    if spec_path.exists():
        try:
            spec = json.loads(spec_path.read_text(encoding="utf-8"))
            slides = [
                {
                    "index": s.get("index", i),
                    "label": s.get("section", ""),
                    "title": (s.get("copy") or {}).get("headline", f"Slide {i+1}"),
                    "body": (s.get("copy") or {}).get("body", ""),
                    "layout": s.get("layout", ""),
                    "grounded": bool((s.get("visual") or {}).get("grounded")),
                    "notes": s.get("speaker_notes", ""),
                }
                for i, s in enumerate(spec.get("slides", []))
            ]
            return {
                "slides": slides,
                "total": len(slides),
                "title": spec.get("title", ""),
                "quality_score": spec.get("quality_score"),
                "grounded": spec.get("grounded", False),
            }
        except Exception:
            pass  # fall through to pptx parsing

    path = Path(settings.upload_dir) / f"{job_id}.pptx"
    if not path.exists():
        raise HTTPException(404, "Deck not found or still generating")

    from pptx import Presentation
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = [sh.text.strip() for sh in slide.shapes if sh.has_text_frame and sh.text.strip()]
        # filter out short page-number tokens (e.g. "01 / 12")
        content = [t for t in texts if not (len(t) < 10 and "/" in t)]
        label   = content[0] if len(content) > 0 else ""
        title   = content[1] if len(content) > 1 else content[0] if content else f"Slide {i+1}"
        body    = content[2] if len(content) > 2 else ""
        slides.append({"index": i, "label": label, "title": title, "body": body})

    return {"slides": slides, "total": len(slides)}
